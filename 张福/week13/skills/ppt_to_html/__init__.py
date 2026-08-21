"""PPT → PNG + HTML 浏览器 技能"""
import io
import os
import re
import shutil
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

from skills.base import Skill

SLIDE_W = 1920
SLIDE_H = 1080
TITLE_COLOR = "#1a237e"
BODY_COLOR = "#333333"
BG_COLOR = "#ffffff"
ACCENT_COLOR = "#e8eaf6"

PROJECT_ROOT = Path(__file__).parent.parent.parent
FILES_DIR = PROJECT_ROOT / "files"

TOOL_DEFINITION = {
    "type": "function",
    "function": {
        "name": "ppt_to_html",
        "description": "将指定的一个或多个 PPT/PPTX 文件逐页转为 PNG 图片，并生成左右结构的 HTML 浏览页面，支持折叠、键盘上下切换。用户上传 ppt 文件后，使用此工具处理。",
        "parameters": {
            "type": "object",
            "properties": {
                "ppt_files": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "要处理的 PPT/PPTX 文件绝对路径列表",
                },
                "session_id": {
                    "type": "string",
                    "description": "当前会话 ID，用于生成唯一输出目录",
                },
            },
            "required": ["ppt_files"],
        },
    },
}

SKILL_DESCRIPTION = """
## PPT 拆解技能
当用户上传了 .ppt/.pptx 文件并要求拆解为 HTML 浏览页面时，使用 `ppt_to_html` 工具。
该工具会：
1. 读取指定的 PPT 文件，逐页渲染为 1920×1080 PNG 图片
2. 每个 PPT 创建独立子目录存放图片
3. 在 files/ppts/ 下生成左右结构 HTML 页面，支持折叠/展开、键盘 ↑↓ 切换
"""


def _num_key(name: str) -> int:
    m = re.match(r"(\d+)", name)
    return int(m.group(1)) if m else 9999


def _sanitize(name: str) -> str:
    name = re.sub(r"[^\w\u4e00-\u9fff\-]", "_", name)
    return name.strip("_") or "ppt"


def _font(size: int) -> ImageFont.FreeTypeFont:
    for p in [
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
        "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]:
        if os.path.exists(p):
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


def _render_slide(slide, slide_idx: int, ppt_name: str) -> Image.Image:
    img = Image.new("RGB", (SLIDE_W, SLIDE_H), BG_COLOR)
    draw = ImageDraw.Draw(img)
    tf = _font(48)
    hf = _font(36)
    bf = _font(28)
    sf = _font(22)

    y = 30
    title_text = None
    bodies = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                if "title" in str(shape.name).lower():
                    title_text = text
                elif para.level == 0 and len(text) < 60:
                    bodies.append(("heading", text))
                else:
                    bodies.append(("body", text))
        if shape.shape_type == 13:
            try:
                blob = shape.image.blob
                ext = shape.image.content_type.split("/")[-1]
                if ext not in ("png", "jpeg", "jpg"):
                    continue
                pil = Image.open(io.BytesIO(blob))
                pil.thumbnail((400, 300), Image.LANCZOS)
                left = int(SLIDE_W * 0.55)
                top = y + 10
                img.paste(pil, (left, top), pil if pil.mode == "RGBA" else None)
            except Exception:
                pass

    if title_text:
        draw.rectangle([(0, 20), (SLIDE_W, 20 + 70)], fill=ACCENT_COLOR)
        draw.text((40, 30), title_text, fill=TITLE_COLOR, font=tf)
        y = 110
    else:
        draw.text((40, 30), f"{ppt_name} - 第 {slide_idx+1} 页", fill=TITLE_COLOR, font=tf)
        y = 100

    y += 10
    for level, text in bodies:
        if y > SLIDE_H - 60:
            break
        if level == "heading":
            draw.text((50, y), f"▪ {text}", fill=TITLE_COLOR, font=hf)
            y += 50
        else:
            lines = []
            cur = ""
            for c in text:
                if draw.textbbox((0, 0), cur + c, font=bf)[2] > SLIDE_W - 120:
                    lines.append(cur)
                    cur = c
                else:
                    cur += c
            if cur:
                lines.append(cur)
            for line in lines[:8]:
                if y > SLIDE_H - 60:
                    break
                draw.text((70, y), line, fill=BODY_COLOR, font=bf)
                y += 38

    pn = f"{slide_idx + 1}"
    bb = draw.textbbox((0, 0), pn, font=sf)
    draw.text(((SLIDE_W - bb[2]) // 2, SLIDE_H - 40), pn, fill="#999999", font=sf)
    draw.line([(0, SLIDE_H - 55), (SLIDE_W, SLIDE_H - 55)], fill="#eeeeee", width=1)
    return img


def _build_html(slides: List[tuple], ppt_dirs: List[Path]) -> str:
    parts = [
        "<!DOCTYPE html><html lang='zh-CN'><head>",
        "<meta charset='UTF-8'><meta name='viewport' content='width=device-width, initial-scale=1.0'>",
        "<title>PPT 课件</title><style>",
        "*{box-sizing:border-box;margin:0;padding:0}",
        "html,body{height:100%;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI','Noto Sans SC',sans-serif;background:#1a1a2e;color:#e0e0e0}",
        ".layout{display:flex;height:100vh}",
        ".sidebar{width:320px;min-width:320px;background:#16213e;overflow-y:auto;border-right:2px solid #0f3460}",
        ".sidebar-header{padding:18px 16px;background:#0f3460;position:sticky;top:0;z-index:10}",
        ".sidebar-header h1{font-size:17px;color:#e94560}",
        ".sidebar-header p{font-size:12px;color:#8899aa;margin-top:3px}",
        ".group-title{padding:10px 16px;font-size:14px;font-weight:700;color:#e94560;background:#1a2744;cursor:pointer;display:flex;justify-content:space-between;align-items:center;user-select:none}",
        ".group-title:hover{background:#233256}",
        ".group-title .arrow{transition:transform .2s;font-size:12px}",
        ".group-title.collapsed .arrow{transform:rotate(-90deg)}",
        ".slide-link{display:block;padding:7px 16px 7px 28px;font-size:13px;color:#aabbcc;text-decoration:none;cursor:pointer;border-left:3px solid transparent;transition:all .12s}",
        ".slide-link:hover{background:#1a2744;color:#fff}",
        ".slide-link.active{background:#0f3460;color:#fff;border-left-color:#e94560;font-weight:600}",
        ".slide-group.collapsed .slide-link{display:none}",
        ".main{flex:1;display:flex;align-items:center;justify-content:center;background:#1a1a2e;padding:20px;overflow:hidden}",
        ".main img{max-width:100%;max-height:100%;object-fit:contain;border-radius:8px;box-shadow:0 4px 20px rgba(0,0,0,.5)}",
        ".placeholder{color:#555;font-size:16px}",
        ".nav-hint{position:fixed;bottom:16px;right:16px;background:rgba(15,52,96,.8);color:#8899aa;padding:8px 14px;border-radius:8px;font-size:12px;z-index:100}",
        "::-webkit-scrollbar{width:6px}",
        "::-webkit-scrollbar-track{background:#16213e}",
        "::-webkit-scrollbar-thumb{background:#0f3460;border-radius:3px}",
        "</style></head><body>",
    ]
    parts.append("<div class='layout'>")
    parts.append("<div class='sidebar'><div class='sidebar-header'>")
    parts.append(f"<h1>📖 PPT 课件</h1><p>{len(ppt_dirs)} 个课件 · {len(slides)} 页</p></div>")

    current = ""
    gi = 0
    for idx, (ppt, img) in enumerate(slides):
        if ppt != current:
            if current:
                parts.append("</div>")
            current = ppt
            gi += 1
            cls = " collapsed" if gi > 1 else ""
            parts.append(f"<div class='group-title{cls}' onclick='toggleGroup(this)'><span>📂 {ppt}</span><span class='arrow'>▼</span></div>")
            parts.append(f"<div class='slide-group{cls}'>")
        num = img.replace("slide_", "").replace(".png", "")
        parts.append(f"<a class='slide-link' href='#' onclick='show({idx});return false'>{num}</a>")
    if slides:
        parts.append("</div>")
    parts.append("</div>")

    paths = [f"'{s[0]}/{s[1]}'" for s in slides]
    parts.append("<div class='main'><div class='placeholder' id='viewer'>← 点击左侧选择页面</div></div>")
    parts.append("<div class='nav-hint'>↑↓ 切换 · 点击课件名折叠</div></div>")
    parts.append(f"<script>const slides=[{','.join(paths)}];let cur=-1;")
    parts.append("function show(i){cur=i;document.querySelectorAll('.slide-link').forEach(e=>e.classList.remove('active'));const a=document.querySelectorAll('.slide-link');if(a[i]){a[i].classList.add('active');a[i].scrollIntoView({block:'nearest',behavior:'smooth'})}document.getElementById('viewer').innerHTML=`<img src='${slides[i]}'>`}")
    parts.append("function toggleGroup(e){e.classList.toggle('collapsed');const b=e.nextElementSibling;if(b&&b.classList.contains('slide-group'))b.classList.toggle('collapsed')}")
    parts.append("document.addEventListener('keydown',function(e){if(e.key==='ArrowDown'){if(cur<0)show(0);else if(cur<slides.length-1)show(cur+1)}if(e.key==='ArrowUp'){if(cur<0)show(0);else if(cur>0)show(cur-1)}})")
    parts.append("</script></body></html>")
    return "\n".join(parts)


def ppt_to_html(ppt_files: List[str], session_id: str = "", output_dir: str = "") -> Dict[str, Any]:
    if not ppt_files:
        return {"success": False, "error": "未提供 PPT 文件路径"}

    if output_dir:
        dst = Path(output_dir).expanduser().resolve()
    else:
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        tag = f"{session_id}_{ts}" if session_id else ts
        dst = FILES_DIR / "ppts" / tag

    if dst.exists():
        shutil.rmtree(dst)
    dst.mkdir(parents=True)

    valid_files = []
    for p in ppt_files:
        pp = Path(p).expanduser().resolve()
        if not pp.exists():
            return {"success": False, "error": f"文件不存在: {pp}"}
        if pp.suffix.lower() not in (".pptx", ".ppt"):
            return {"success": False, "error": f"不支持的文件格式: {pp.name}"}
        valid_files.append(pp)

    all_slides = []
    ppt_names = []
    for pptx in valid_files:
        dir_name = _sanitize(pptx.stem)
        sub = dst / dir_name
        sub.mkdir(parents=True)
        ppt_names.append(pptx.stem)

        try:
            prs = Presentation(str(pptx))
        except Exception as e:
            return {"success": False, "error": f"无法打开 {pptx.name}: {e}"}

        for i, slide in enumerate(prs.slides):
            try:
                pil_img = _render_slide(slide, i, pptx.stem)
                fname = f"slide_{i+1:03d}.png"
                pil_img.save(sub / fname, "PNG")
                all_slides.append((dir_name, fname))
            except Exception as e:
                all_slides.append((dir_name, f"slide_{i+1:03d}.png"))

    html = _build_html(all_slides, ppt_names)
    (dst / "index.html").write_text(html, encoding="utf-8")

    return {
        "success": True,
        "message": f"完成！{len(valid_files)} 个课件，{len(all_slides)} 页幻灯片",
        "output_dir": str(dst),
        "html_file": f"{dst}/index.html",
        "total_ppts": len(valid_files),
        "total_slides": len(all_slides),
    }


class PptToHtmlSkill(Skill):
    name = "ppt_to_html"
    description = SKILL_DESCRIPTION
    tools = [TOOL_DEFINITION]

    def execute(self, tool_name: str, arguments: Dict[str, Any]) -> Dict[str, Any]:
        if tool_name != "ppt_to_html":
            return {"success": False, "error": f"未知工具: {tool_name}"}
        return ppt_to_html(**arguments)


skill = PptToHtmlSkill()
